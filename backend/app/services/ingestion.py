import re
from pathlib import Path

from pypdf import PdfReader
from pptx import Presentation

from ..config import settings


def extract_text(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(file_path)
    elif ext == ".pptx":
        return _extract_pptx(file_path)
    elif ext in (".txt", ".md"):
        return Path(file_path).read_text(encoding="utf-8", errors="ignore")
    raise ValueError(f"Unsupported file type: {ext}")


def _extract_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)
    pages = []
    extracted_pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        extracted_pages.append(text)
        pages.append(f"[page {i + 1}]\n{text}")
    if not any(text.strip() for text in extracted_pages):
        raise ValueError(
            "No extractable text was found in this PDF. It may be a scanned image; "
            "run OCR on it before uploading."
        )
    return "\n\n".join(pages)


def _extract_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        parts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
        text = "\n".join(parts)
        if notes.strip():
            text += f"\n[speaker notes] {notes}"
        slides.append(f"[slide {i + 1}]\n{text}")
    return "\n\n".join(slides)


def chunk_text(text: str, chunk_size: int | None = None, overlap: int | None = None) -> list[str]:
    """Pack paragraphs into ~chunk_size windows, carrying a small overlap
    into the next chunk so concepts near a boundary aren't split away from
    their context."""
    chunk_size = chunk_size or settings.chunk_size
    overlap = overlap or settings.chunk_overlap

    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    chunks: list[str] = []
    current = ""
    for para in paragraphs:
        if len(current) + len(para) + 1 <= chunk_size:
            current = f"{current}\n{para}".strip()
        else:
            if current:
                chunks.append(current)
            tail = current[-overlap:] if current else ""
            current = f"{tail}\n{para}".strip()
    if current:
        chunks.append(current)
    return chunks
